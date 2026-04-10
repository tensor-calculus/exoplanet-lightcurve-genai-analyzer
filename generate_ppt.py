"""
Generate a PowerPoint presentation for the Exoplanet Lightcurve Analyzer project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color Palette ---
DEEP_NAVY = RGBColor(0x0F, 0x0C, 0x29)
DARK_BLUE = RGBColor(0x1A, 0x1A, 0x2E)
MID_BLUE = RGBColor(0x16, 0x21, 0x3E)
ACCENT_PURPLE = RGBColor(0x66, 0x7E, 0xEA)
ACCENT_VIOLET = RGBColor(0x76, 0x4B, 0xA2)
ACCENT_PINK = RGBColor(0xF0, 0x93, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xD0, 0xD8, 0xF0)
MUTED_TEXT = RGBColor(0x88, 0x92, 0xA4)
SOFT_CYAN = RGBColor(0x4A, 0x90, 0xD9)
SUCCESS_GREEN = RGBColor(0x00, 0xC9, 0xA7)
WARM_ORANGE = RGBColor(0xFF, 0xA7, 0x51)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_shape_alpha(shape, alpha):
    """Set alpha transparency on a shape's fill using raw XML manipulation."""
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    # Access the shape's XML element
    spPr = shape._element.find('.//a:solidFill/a:srgbClr', nsmap)
    if spPr is not None:
        alpha_elem = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
        alpha_elem.set('val', str(int(alpha * 1000)))


def add_shape(slide, left, top, width, height, fill_color, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                        color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                        font_name="Calibri", line_spacing=1.5):
    """Add a text box with multiple lines (paragraphs)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16,
                       color=LIGHT_TEXT, font_name="Calibri", bullet_color=ACCENT_PURPLE):
    """Add bullet-pointed text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_PURPLE):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ========== SLIDE 1: TITLE SLIDE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, DEEP_NAVY)

    # Decorative shapes
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)
    add_shape(slide, Inches(8.5), Inches(0), Inches(4.833), Inches(7.5), MID_BLUE, alpha=30)

    # Emoji planet
    add_text_box(slide, Inches(1.5), Inches(0.8), Inches(2), Inches(1.5),
                 "🪐", font_size=80, alignment=PP_ALIGN.LEFT)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.2), Inches(7), Inches(1),
                 "Exoplanet Lightcurve Analyzer", font_size=44, color=WHITE,
                 bold=True, font_name="Calibri")

    # Accent line
    add_accent_line(slide, Inches(1.5), Inches(3.4), Inches(2.5))

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(3.7), Inches(7), Inches(0.8),
                 "AI-Powered Visual Similarity Search for Exoplanet Transit Light Curves",
                 font_size=20, color=LIGHT_TEXT)

    # Description
    add_text_box(slide, Inches(1.5), Inches(4.6), Inches(6.5), Inches(1.5),
                 "A Generative AI application that uses CLIP embeddings and RAG-based chat "
                 "to identify and analyze exoplanets through their transit light curve signatures.",
                 font_size=14, color=MUTED_TEXT)

    # Tech tags at bottom
    tags = ["CLIP  •  ChromaDB  •  Streamlit  •  GPT-3.5  •  NASA Exoplanet Archive"]
    add_text_box(slide, Inches(1.5), Inches(6.3), Inches(8), Inches(0.5),
                 tags[0], font_size=12, color=ACCENT_PURPLE, font_name="Calibri")

    # Right side - sample lightcurve image
    sample_img = "data/images/Kepler-45_b.png"
    if os.path.exists(sample_img):
        slide.shapes.add_picture(sample_img, Inches(9), Inches(1.5), Inches(3.5), Inches(3.5))
        add_text_box(slide, Inches(9), Inches(5.1), Inches(3.5), Inches(0.4),
                     "Sample: Kepler-45 b Light Curve", font_size=11, color=MUTED_TEXT,
                     alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 2: PROBLEM STATEMENT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    # Section number
    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "01", font_size=36, color=ACCENT_PURPLE, bold=True)

    add_text_box(slide, Inches(1), Inches(1.0), Inches(5), Inches(0.7),
                 "Problem Statement", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Left column - Problem
    add_rounded_rect(slide, Inches(1), Inches(2.3), Inches(5.5), Inches(4.5), MID_BLUE, alpha=50)
    add_text_box(slide, Inches(1.4), Inches(2.5), Inches(4.8), Inches(0.5),
                 "🔭  The Challenge", font_size=20, color=ACCENT_PINK, bold=True)
    add_bullet_points(slide, Inches(1.4), Inches(3.2), Inches(4.8), Inches(3.5), [
        "Thousands of exoplanets discovered via transit method",
        "Light curve patterns are complex and hard to compare manually",
        "Astronomers need fast similarity matching across large datasets",
        "Traditional keyword search doesn't capture visual features",
        "Lack of accessible AI tools for light curve analysis",
    ], font_size=14)

    # Right column - Opportunity
    add_rounded_rect(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4.5), MID_BLUE, alpha=50)
    add_text_box(slide, Inches(7.4), Inches(2.5), Inches(4.8), Inches(0.5),
                 "💡  The Opportunity", font_size=20, color=SUCCESS_GREEN, bold=True)
    add_bullet_points(slide, Inches(7.4), Inches(3.2), Inches(4.8), Inches(3.5), [
        "Use AI to encode visual light curve features into embeddings",
        "Enable image-based similarity search across planet databases",
        "Let AI provide grounded explanations of matched results",
        "Make exoplanet analysis accessible to researchers & students",
        "Bridge computer vision (CLIP) with astronomy",
    ], font_size=14)

    # ========== SLIDE 3: TECH STACK & ARCHITECTURE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "02", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(6), Inches(0.7),
                 "Architecture & Technology Stack", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Tech cards
    techs = [
        ("🧠", "CLIP (ViT-B/32)", "OpenAI's vision-language model.\nConverts light curve images\ninto 512-dim vector embeddings.", ACCENT_PURPLE),
        ("🗄️", "ChromaDB", "Persistent vector database.\nStores embeddings with metadata.\nCosine similarity search.", SOFT_CYAN),
        ("🌐", "Streamlit", "Interactive web application.\nReal-time image upload,\nresults display & chat UI.", SUCCESS_GREEN),
        ("🤖", "GPT-3.5 Turbo", "OpenAI language model.\nRAG-grounded responses about\nmatched exoplanet data.", ACCENT_PINK),
        ("🔭", "NASA Archive", "Real data from NASA\nExoplanet Archive TAP API.\n55 confirmed transit planets.", WARM_ORANGE),
    ]

    card_width = Inches(2.1)
    card_height = Inches(3.2)
    start_x = Inches(0.7)
    spacing = Inches(0.3)

    for i, (icon, title, desc, accent) in enumerate(techs):
        x = start_x + i * (card_width + spacing)
        y = Inches(2.4)

        card = add_rounded_rect(slide, x, y, card_width, card_height, DARK_BLUE, alpha=80)

        # Accent top bar
        add_shape(slide, x, y, card_width, Pt(4), accent)

        # Icon
        add_text_box(slide, x, y + Inches(0.15), card_width, Inches(0.7),
                     icon, font_size=36, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.15), y + Inches(0.85), card_width - Inches(0.3), Inches(0.4),
                     title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        desc_lines = desc.split('\n')
        add_multiline_text(slide, x + Inches(0.15), y + Inches(1.4),
                           card_width - Inches(0.3), Inches(1.5),
                           desc_lines, font_size=11, color=MUTED_TEXT,
                           alignment=PP_ALIGN.CENTER)

    # Pipeline flow at bottom
    add_rounded_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.0), MID_BLUE, alpha=40)
    add_text_box(slide, Inches(1), Inches(6.15), Inches(11.5), Inches(0.7),
                 "Pipeline:   NASA API  →  fetch_data.py  →  Light Curve Images + Metadata  →  "
                 "build_index.py  →  CLIP Embeddings  →  ChromaDB  →  app.py (Streamlit + RAG Chat)",
                 font_size=13, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 4: DATA PIPELINE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "03", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(6), Inches(0.7),
                 "Data Pipeline — fetch_data.py", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Step boxes
    steps = [
        ("Step 1: Query NASA", "Queries the NASA Exoplanet\nArchive TAP API for confirmed\ntransiting planets with known\nperiod, temperature, and radius.",
         "🛰️", ACCENT_PURPLE),
        ("Step 2: Generate Curves", "Creates synthetic folded transit\nlight curves using real planet\nparameters (Rp/Rs ratio,\norbital period, noise modeling).",
         "📊", SOFT_CYAN),
        ("Step 3: Classify", "Classifies each planet by radius:\nTerrestrial, Super-Earth,\nSub-Neptune, Neptune-like,\nor Gas Giant.",
         "🏷️", SUCCESS_GREEN),
        ("Step 4: Save Data", "Saves 400×400px PNG images\nand JSON metadata files for\neach planet. 55 planets\nprocessed in total.",
         "💾", WARM_ORANGE),
    ]

    for i, (title, desc, icon, accent) in enumerate(steps):
        x = Inches(1) + i * Inches(3)
        y = Inches(2.4)

        add_rounded_rect(slide, x, y, Inches(2.7), Inches(3.5), DARK_BLUE, alpha=70)
        add_shape(slide, x, y, Inches(2.7), Pt(4), accent)

        add_text_box(slide, x, y + Inches(0.2), Inches(2.7), Inches(0.6),
                     icon, font_size=32, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(2.3), Inches(0.4),
                     title, font_size=15, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

        desc_lines = desc.split('\n')
        add_multiline_text(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.3), Inches(2),
                           desc_lines, font_size=12, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)

    # Sample images row
    sample_images = ["data/images/CoRoT-11_b.png", "data/images/Kepler-18_c.png",
                     "data/images/WASP-23_b.png", "data/images/NGTS-10_b.png"]
    img_y = Inches(6.2)
    for i, img_path in enumerate(sample_images):
        if os.path.exists(img_path):
            x = Inches(1.5) + i * Inches(2.8)
            slide.shapes.add_picture(img_path, x, img_y, Inches(1.0), Inches(1.0))
            name = os.path.splitext(os.path.basename(img_path))[0].replace("_", " ")
            add_text_box(slide, x - Inches(0.3), img_y + Inches(1.0), Inches(1.6), Inches(0.3),
                         name, font_size=9, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 5: EMBEDDING & INDEXING ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "04", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(8), Inches(0.7),
                 "Embedding & Vector Indexing — build_index.py", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Left: CLIP embedding explanation
    add_rounded_rect(slide, Inches(1), Inches(2.3), Inches(5.8), Inches(4.7), MID_BLUE, alpha=40)
    add_text_box(slide, Inches(1.4), Inches(2.5), Inches(5), Inches(0.5),
                 "🧠  CLIP Visual Embeddings", font_size=20, color=ACCENT_PURPLE, bold=True)

    add_bullet_points(slide, Inches(1.4), Inches(3.2), Inches(5), Inches(3.5), [
        "Each light curve image → 512-dimensional vector",
        "CLIP model: openai/clip-vit-base-patch32",
        "Embeddings are L2-normalized for cosine similarity",
        "Batch processing (16 images per batch) for efficiency",
        "Captures visual features: transit depth, shape, noise",
        "Similar light curves → nearby vectors in embedding space",
    ], font_size=14)

    # Right: ChromaDB details
    add_rounded_rect(slide, Inches(7.3), Inches(2.3), Inches(5.3), Inches(4.7), MID_BLUE, alpha=40)
    add_text_box(slide, Inches(7.7), Inches(2.5), Inches(4.5), Inches(0.5),
                 "🗄️  ChromaDB Vector Store", font_size=20, color=SOFT_CYAN, bold=True)

    add_bullet_points(slide, Inches(7.7), Inches(3.2), Inches(4.5), Inches(3.5), [
        "Persistent local vector database",
        "Collection: 'exoplanets' with cosine distance",
        "Stores embeddings + rich metadata per planet",
        "Metadata: name, period, temp, radius, type, star",
        "HNSW index for fast approximate nearest neighbor",
        "55 planets indexed with full metadata",
    ], font_size=14)

    # ========== SLIDE 6: APPLICATION FEATURES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "05", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(8), Inches(0.7),
                 "Application Features — app.py", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Feature cards - 3 across
    features = [
        ("📤", "Image Upload & Search", [
            "Upload any light curve image (PNG/JPG)",
            "CLIP encodes it to a 512-dim vector",
            "Queries ChromaDB for top-N closest matches",
            "Adjustable number of results (1–10)",
        ], ACCENT_PURPLE),
        ("🪐", "Planet Result Cards", [
            "Beautiful styled cards for each match",
            "Shows similarity score as percentage",
            "Displays: period, temp, radius, mass, star",
            "Planet type classification badge",
        ], SOFT_CYAN),
        ("💬", "RAG-Grounded AI Chat", [
            "Chat with GPT-3.5 about matched planets",
            "System prompt grounded in actual data",
            "References real planet names & values",
            "Educational, accurate, engaging responses",
        ], ACCENT_PINK),
    ]

    for i, (icon, title, items, accent) in enumerate(features):
        x = Inches(0.7) + i * Inches(4.2)
        y = Inches(2.4)

        add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.3), DARK_BLUE, alpha=70)
        add_shape(slide, x, y, Inches(3.8), Pt(4), accent)

        add_text_box(slide, x, y + Inches(0.2), Inches(3.8), Inches(0.6),
                     icon, font_size=36, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(0.4),
                     title, font_size=16, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

        add_bullet_points(slide, x + Inches(0.3), y + Inches(1.5), Inches(3.2), Inches(2.5),
                          items, font_size=12, color=LIGHT_TEXT)

    # ========== SLIDE 7: RAG PIPELINE DEEP DIVE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "06", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(8), Inches(0.7),
                 "RAG Pipeline — Retrieval Augmented Generation", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Flow diagram as horizontal steps
    rag_steps = [
        ("1️⃣", "User uploads\nlight curve image", ACCENT_PURPLE),
        ("2️⃣", "CLIP generates\n512-dim embedding", SOFT_CYAN),
        ("3️⃣", "ChromaDB finds\ntop-N matches", SUCCESS_GREEN),
        ("4️⃣", "Metadata injected\ninto system prompt", WARM_ORANGE),
        ("5️⃣", "GPT-3.5 answers\ngrounded in data", ACCENT_PINK),
    ]

    for i, (num, desc, accent) in enumerate(rag_steps):
        x = Inches(0.5) + i * Inches(2.5)
        y = Inches(2.5)

        add_rounded_rect(slide, x, y, Inches(2.2), Inches(2.2), DARK_BLUE, alpha=70)
        add_shape(slide, x, y, Inches(2.2), Pt(4), accent)

        add_text_box(slide, x, y + Inches(0.15), Inches(2.2), Inches(0.5),
                     num, font_size=28, alignment=PP_ALIGN.CENTER)

        desc_lines = desc.split('\n')
        add_multiline_text(slide, x + Inches(0.15), y + Inches(0.7),
                           Inches(1.9), Inches(1.3),
                           desc_lines, font_size=13, color=LIGHT_TEXT,
                           alignment=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < len(rag_steps) - 1:
            add_text_box(slide, x + Inches(2.15), y + Inches(0.7), Inches(0.4), Inches(0.5),
                         "→", font_size=24, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

    # Benefits section
    add_rounded_rect(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(1.8), MID_BLUE, alpha=40)
    add_text_box(slide, Inches(1.4), Inches(5.4), Inches(3), Inches(0.4),
                 "Key Benefits of RAG Approach", font_size=16, color=SUCCESS_GREEN, bold=True)

    benefits = [
        "✅  Responses are grounded in actual exoplanet data — no hallucination about planet properties",
        "✅  Users can ask follow-up questions about specific matched planets with accurate context",
        "✅  System clearly indicates when going beyond retrieved data into general knowledge",
    ]
    add_multiline_text(slide, Inches(1.4), Inches(5.9), Inches(10.5), Inches(1.0),
                       benefits, font_size=12, color=LIGHT_TEXT)

    # ========== SLIDE 8: PROJECT STRUCTURE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "07", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(6), Inches(0.7),
                 "Project Structure & Codebase", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # File structure
    files_info = [
        ("fetch_data.py", "Data Acquisition", "Queries NASA Exoplanet Archive API, generates synthetic transit\nlight curves, classifies planets, saves images + metadata.", "231 lines", ACCENT_PURPLE),
        ("build_index.py", "Vector Indexing", "Loads CLIP model, processes images in batches of 16, generates\nnormalized embeddings, stores in ChromaDB with metadata.", "151 lines", SOFT_CYAN),
        ("app.py", "Web Application", "Streamlit UI with image upload, CLIP-based query, styled result\ncards, and RAG chat with GPT-3.5 grounded in matched data.", "379 lines", SUCCESS_GREEN),
        ("requirements.txt", "Dependencies", "Core deps: transformers, torch, chromadb, streamlit, openai,\npillow, matplotlib, requests, numpy.", "11 lines", WARM_ORANGE),
    ]

    for i, (fname, role, desc, lines, accent) in enumerate(files_info):
        y = Inches(2.2) + i * Inches(1.25)
        add_rounded_rect(slide, Inches(1), y, Inches(11.3), Inches(1.1), DARK_BLUE, alpha=60)
        add_shape(slide, Inches(1), y, Pt(4), Inches(1.1), accent)

        add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(2.2), Inches(0.4),
                     fname, font_size=16, color=accent, bold=True)
        add_text_box(slide, Inches(3.5), y + Inches(0.05), Inches(2), Inches(0.4),
                     role, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(10.5), y + Inches(0.05), Inches(1.5), Inches(0.4),
                     lines, font_size=12, color=MUTED_TEXT, alignment=PP_ALIGN.RIGHT)

        desc_lines = desc.split('\n')
        add_multiline_text(slide, Inches(3.5), y + Inches(0.4), Inches(7), Inches(0.7),
                           desc_lines, font_size=11, color=MUTED_TEXT)

    # Data stats
    add_rounded_rect(slide, Inches(1), Inches(7.5) - Inches(1.2), Inches(11.3), Inches(0.9), MID_BLUE, alpha=40)

    stats = [
        ("55", "Exoplanets"),
        ("512", "Embedding Dim"),
        ("~770", "Lines of Code"),
        ("10", "Dependencies"),
    ]
    for i, (val, label) in enumerate(stats):
        x = Inches(1.5) + i * Inches(2.8)
        y_stat = Inches(7.5) - Inches(1.1)
        add_text_box(slide, x, y_stat, Inches(2), Inches(0.4),
                     val, font_size=24, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y_stat + Inches(0.35), Inches(2), Inches(0.3),
                     label, font_size=12, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 9: HOW TO RUN ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "08", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(6), Inches(0.7),
                 "Getting Started", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Setup steps
    run_steps = [
        ("1", "Install Dependencies", "pip install -r requirements.txt", "Install all required Python packages including\nCLIP, ChromaDB, Streamlit, and OpenAI SDK."),
        ("2", "Fetch Exoplanet Data", "python fetch_data.py", "Downloads real exoplanet data from NASA Archive\nand generates synthetic light curve images."),
        ("3", "Build Vector Index", "python build_index.py", "Creates CLIP embeddings for all light curves and\nstores them in ChromaDB vector database."),
        ("4", "Launch Application", "streamlit run app.py", "Opens the web app where you can upload light\ncurves and chat with the AI astronomer."),
    ]

    for i, (num, title, cmd, desc) in enumerate(run_steps):
        x = Inches(1)
        y = Inches(2.3) + i * Inches(1.2)

        add_rounded_rect(slide, x, y, Inches(11.3), Inches(1.05), DARK_BLUE, alpha=60)

        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_PURPLE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)

        add_text_box(slide, x + Inches(1), y + Inches(0.08), Inches(3), Inches(0.4),
                     title, font_size=16, color=WHITE, bold=True)

        # Command in code-style
        add_rounded_rect(slide, x + Inches(1), y + Inches(0.5), Inches(3.5), Inches(0.4), RGBColor(0x0A, 0x0A, 0x1A))
        add_text_box(slide, x + Inches(1.15), y + Inches(0.52), Inches(3.3), Inches(0.35),
                     f"$ {cmd}", font_size=11, color=SUCCESS_GREEN, font_name="Consolas")

        desc_lines = desc.split('\n')
        add_multiline_text(slide, x + Inches(5), y + Inches(0.15), Inches(6), Inches(0.8),
                           desc_lines, font_size=12, color=MUTED_TEXT)

    # ========== SLIDE 10: FUTURE SCOPE & THANK YOU ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(1), Inches(0.6),
                 "09", font_size=36, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(1), Inches(1.0), Inches(6), Inches(0.7),
                 "Future Scope & Conclusion", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2))

    # Future scope
    add_rounded_rect(slide, Inches(1), Inches(2.3), Inches(5.8), Inches(3.5), MID_BLUE, alpha=40)
    add_text_box(slide, Inches(1.4), Inches(2.5), Inches(5), Inches(0.5),
                 "🚀  Future Enhancements", font_size=18, color=ACCENT_PINK, bold=True)

    future_items = [
        "Scale to 500+ exoplanets from expanded NASA queries",
        "Add real light curve data from TESS / Kepler missions",
        "Fine-tune CLIP on actual transit light curve datasets",
        "Multi-modal search: combine image + text queries",
        "Deploy on cloud (Streamlit Cloud / HuggingFace Spaces)",
        "Add export and comparison features for researchers",
    ]
    add_bullet_points(slide, Inches(1.4), Inches(3.1), Inches(5), Inches(2.5),
                      future_items, font_size=13)

    # Conclusion / Thank you
    add_rounded_rect(slide, Inches(7.3), Inches(2.3), Inches(5.3), Inches(3.5), MID_BLUE, alpha=40)
    add_text_box(slide, Inches(7.7), Inches(2.5), Inches(4.5), Inches(0.5),
                 "🎯  Key Takeaways", font_size=18, color=SUCCESS_GREEN, bold=True)

    takeaways = [
        "Successfully combined CLIP + ChromaDB + LLM for astronomy",
        "Image-based similarity search for exoplanet light curves",
        "RAG-grounded chat eliminates AI hallucination",
        "Real NASA data ensures scientific accuracy",
        "Clean, modular, 3-file architecture",
        "Extensible foundation for research use",
    ]
    add_bullet_points(slide, Inches(7.7), Inches(3.1), Inches(4.5), Inches(2.5),
                      takeaways, font_size=13)

    # Thank you banner
    add_rounded_rect(slide, Inches(2.5), Inches(6.2), Inches(8.3), Inches(1.0), ACCENT_PURPLE, alpha=30)
    add_text_box(slide, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.5),
                 "Thank You!", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.5), Inches(6.8), Inches(8.3), Inches(0.4),
                 "Exoplanet Lightcurve Analyzer  •  Powered by GenAI",
                 font_size=14, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

    # --- Save ---
    output_path = "ExoplanetAnalyzer_Presentation.pptx"
    prs.save(output_path)
    print(f"\n[OK] Presentation saved to: {output_path}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
